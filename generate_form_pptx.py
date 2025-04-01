from pptx import Presentation
import qrcode
import os
from io import BytesIO
import sqlite3
import datetime
from kivy.app import App
from kivy.uix.boxlayout import BoxLayout
from kivy.uix.gridlayout import GridLayout
from kivy.uix.label import Label
from kivy.uix.textinput import TextInput
from kivy.uix.button import Button
from kivy.uix.togglebutton import ToggleButton
from kivy.uix.popup import Popup
from kivy.core.window import Window
from kivy.uix.scrollview import ScrollView
from kivy.metrics import dp

def create_database():
    """Создание базы данных, если она не существует"""
    conn = sqlite3.connect('маршрутные_карты.db')
    cursor = conn.cursor()
    cursor.execute('''
    CREATE TABLE IF NOT EXISTS маршрутные_карты (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        Номер_бланка TEXT NOT NULL,
        Учетный_номер TEXT,
        Номер_кластера TEXT,
        Статус TEXT,
        Дата_создания TEXT,
        Путь_к_файлу TEXT
    )
    ''')
    conn.commit()
    conn.close()

def save_to_database(form_number, file_path):
    """Сохранение информации о созданной маршрутной карте в базу данных"""
    try:
        with sqlite3.connect('маршрутные_карты.db') as conn:
            cursor = conn.cursor()
            date_created = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            cursor.execute(
                "INSERT INTO маршрутные_карты (Номер_бланка, Учетный_номер, Номер_кластера, Статус, Дата_создания, Путь_к_файлу) VALUES (?, ?, ?, ?, ?, ?)",
                (form_number, None, None, None, date_created, file_path)
            )
            conn.commit()
    except sqlite3.Error as e:
        print(f"Ошибка при сохранении в базу данных: {e}")
        raise

def check_duplicate_form_number(form_number):
    """Проверка существования бланка с таким номером в базе данных"""
    conn = sqlite3.connect('маршрутные_карты.db')
    cursor = conn.cursor()
    cursor.execute("SELECT COUNT(*) FROM маршрутные_карты WHERE Номер_бланка = ?", (form_number,))
    count = cursor.fetchone()[0]
    conn.close()
    return count > 0

def generate_form_with_qr(template_path, output_path, form_number):
    try:
        # Проверяем, существует ли уже бланк с таким номером
        if check_duplicate_form_number(form_number):
            raise ValueError(f"Бланк с номером {form_number} уже существует в базе данных")
        
        # Открываем шаблон презентации
        prs = Presentation(template_path)
        
        # Создаем QR-код с номером формы
        qr = qrcode.QRCode(
            version=1,
            error_correction=qrcode.constants.ERROR_CORRECT_L,
            box_size=10,
            border=1,
        )
        qr.add_data(form_number)
        qr.make(fit=True)
        qr_image = qr.make_image(fill_color="black", back_color="white")
        
        # Сохраняем QR-код во временный буфер
        image_stream = BytesIO()
        qr_image.save(image_stream, format='PNG')
        image_stream.seek(0)
        
        # Добавляем QR-код на первый слайд
        slide = prs.slides[0]
        
        # Ищем текст "МАРШРУТНАЯ КАРТА"
        target_text = "МАРШРУТНАЯ КАРТА"
        text_shape = None
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and target_text in shape.text:
                text_shape = shape
                break
        
        # Задаем размеры QR-кода
        qr_width = 500000  # ~0.5 см
        qr_height = 500000
        
        if text_shape:
            # Располагаем QR-код слева от текста на том же уровне
            left = max(300000, text_shape.left - qr_width - 200000)  # отступ от текста ~0.2 см, но не меньше 300000
            
            # Выравниваем по вертикали с текстом
            top = text_shape.top + (text_shape.height - qr_height) / 2
            
            # Добавляем QR-код на слайд
            slide.shapes.add_picture(
                image_stream,
                left,
                top,
                width=qr_width,
                height=qr_height
            )
            
            # Добавляем номер в правый верхний угол слайда
            slide_width = prs.slide_width
            number_shape = slide.shapes.add_textbox(
                slide_width - 1500000,  # отступ от правого края
                200000,  # отступ от верхнего края
                1200000,  # ширина текстового поля
                300000  # высота текстового поля
            )
            text_frame = number_shape.text_frame
            p = text_frame.paragraphs[0]
            p.text = f"№ {form_number}"
            p.alignment = 2  # выравнивание по правому краю
            
        else:
            # Если текст не найден, используем позицию по умолчанию
            left = 300000  # отступ от левого края
            top = 300000  # отступ от верхнего края
            
            # Добавляем QR-код на слайд
            slide.shapes.add_picture(
                image_stream,
                left,
                top,
                width=qr_width,
                height=qr_height
            )
            
            # Добавляем номер в правый верхний угол
            slide_width = prs.slide_width
            number_shape = slide.shapes.add_textbox(
                slide_width - 1500000,  # отступ от правого края
                200000,  # отступ от верхнего края
                1200000,  # ширина текстового поля
                300000  # высота текстового поля
            )
            text_frame = number_shape.text_frame
            p = text_frame.paragraphs[0]
            p.text = f"№ {form_number}"
            p.alignment = 2  # выравнивание по правому краю
        
        # Сохраняем результат
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        
        # Сохраняем информацию в базу данных
        save_to_database(form_number, output_path)
    except Exception as e:
        print(f"Ошибка при генерации формы: {e}")
        raise

def generate_multiple_forms(template_path, start_number, count):
    """Генерация нескольких маршрутных карт"""
    # Создаем папку для сохранения файлов, если она не существует
    output_dir = "Маршрутные_карты"
    os.makedirs(output_dir, exist_ok=True)
    
    # Проверяем все номера на дубликаты перед генерацией
    start_num = int(start_number)
    duplicates = []
    
    for i in range(count):
        current_num = start_num + i
        form_number = f"{current_num:06d}"
        if check_duplicate_form_number(form_number):
            duplicates.append(form_number)
    
    if duplicates:
        return 0, [f"Следующие номера бланков уже существуют в базе данных: {', '.join(duplicates)}"]
    
    # Если дубликатов нет, продолжаем генерацию
    errors = []
    success_count = 0
    
    for i in range(count):
        current_num = start_num + i
        # Форматируем номер с ведущими нулями (6 цифр)
        form_number = f"{current_num:06d}"
        output_path = os.path.join(output_dir, f"маршрутная_карта_{form_number}.pptx")
        
        try:
            generate_form_with_qr(template_path, output_path, form_number)
            success_count += 1
        except Exception as e:
            errors.append(f"Ошибка при создании файла {output_path}: {str(e)}")
    
    return success_count, errors

class InfoPopup(Popup):
    def __init__(self, title, message, **kwargs):
        super(InfoPopup, self).__init__(**kwargs)
        self.title = title
        self.size_hint = (0.8, 0.4)
        
        content = BoxLayout(orientation='vertical', padding=dp(10), spacing=dp(10))
        
        # Создаем прокручиваемую область для длинных сообщений
        scroll_view = ScrollView(size_hint=(1, 0.8))
        message_label = Label(text=message, size_hint_y=None, halign='left', valign='top')
        message_label.bind(width=lambda *x: setattr(message_label, 'text_size', (message_label.width, None)))
        message_label.bind(texture_size=lambda *x: setattr(message_label, 'height', message_label.texture_size[1]))
        scroll_view.add_widget(message_label)
        
        # Кнопка закрытия
        close_button = Button(text="Закрыть", size_hint=(1, 0.2))
        close_button.bind(on_press=self.dismiss)
        
        content.add_widget(scroll_view)
        content.add_widget(close_button)
        
        self.content = content

class FormGeneratorUI(BoxLayout):
    def __init__(self, **kwargs):
        super(FormGeneratorUI, self).__init__(**kwargs)
        self.orientation = 'vertical'
        self.padding = dp(20)
        self.spacing = dp(10)
        
        # Создаем базу данных при запуске
        create_database()
        
        # Заголовок
        self.add_widget(Label(
            text="Генератор маршрутных карт",
            font_size=dp(24),
            size_hint=(1, 0.1)
        ))
        
        # Режим генерации
        mode_box = BoxLayout(orientation='vertical', size_hint=(1, 0.2))
        mode_label = Label(text="Режим генерации:", halign='left', size_hint=(1, 0.3))
        mode_label.bind(size=lambda *x: setattr(mode_label, 'text_size', (mode_label.width, None)))
        
        mode_buttons = BoxLayout(size_hint=(1, 0.7))
        self.single_mode = ToggleButton(text="Один бланк", group="mode", state="down")
        self.multiple_mode = ToggleButton(text="Несколько бланков", group="mode")
        
        self.single_mode.bind(on_press=self.toggle_mode)
        self.multiple_mode.bind(on_press=self.toggle_mode)
        
        mode_buttons.add_widget(self.single_mode)
        mode_buttons.add_widget(self.multiple_mode)
        
        mode_box.add_widget(mode_label)
        mode_box.add_widget(mode_buttons)
        self.add_widget(mode_box)
        
        # Параметры
        params_grid = GridLayout(cols=2, size_hint=(1, 0.3), spacing=dp(10))
        
        # Номер бланка
        params_grid.add_widget(Label(text="Номер бланка:", halign='left'))
        self.form_number = TextInput(text="000001", multiline=False, input_filter='int')
        params_grid.add_widget(self.form_number)
        
        # Количество бланков
        params_grid.add_widget(Label(text="Количество бланков:", halign='left'))
        self.count_input = TextInput(text="50", multiline=False, input_filter='int', disabled=True)
        params_grid.add_widget(self.count_input)
        
        self.add_widget(params_grid)
        
        # Кнопки
        buttons_box = BoxLayout(size_hint=(1, 0.1), spacing=dp(10))
        
        generate_button = Button(text="Создать")
        generate_button.bind(on_press=self.generate_forms)
        
        exit_button = Button(text="Выход")
        exit_button.bind(on_press=lambda x: App.get_running_app().stop())
        
        buttons_box.add_widget(generate_button)
        buttons_box.add_widget(exit_button)
        
        self.add_widget(buttons_box)
        
        # Статус
        self.status_label = Label(
            text="",
            size_hint=(1, 0.1),
            color=(0, 0.7, 1, 1)
        )
        self.add_widget(self.status_label)
    
    def toggle_mode(self, instance):
        """Переключение между режимами генерации одного или нескольких бланков"""
        if self.single_mode.state == "down":
            self.count_input.disabled = True
        else:
            self.count_input.disabled = False
    
    def show_popup(self, title, message):
        """Показать всплывающее окно с сообщением"""
        popup = InfoPopup(title=title, message=message)
        popup.open()
    
    def generate_forms(self, instance):
        """Генерация маршрутных карт"""
        template_path = "ШАБЛОН.pptx"
        
        # Проверяем наличие шаблона
        if not os.path.exists(template_path):
            self.show_popup("Ошибка", f"Файл шаблона '{template_path}' не найден!")
            return
        
        # Проверяем корректность номера бланка
        form_number = self.form_number.text.strip()
        if not form_number.isdigit():
            self.show_popup("Ошибка", "Номер бланка должен содержать только цифры!")
            return
        
        # Создаем папку для сохранения файлов, если она не существует
        output_dir = "Маршрутные_карты"
        os.makedirs(output_dir, exist_ok=True)
        
        if self.single_mode.state == "down":
            # Генерация одного бланка
            form_number_formatted = f"{int(form_number):06d}"
            output_path = os.path.join(output_dir, f"маршрутная_карта_{form_number_formatted}.pptx")
            
            try:
                generate_form_with_qr(template_path, output_path, form_number_formatted)
                self.status_label.text = f"Файл успешно создан: {output_path}"
                self.show_popup("Успех", f"Маршрутная карта успешно создана!\nФайл: {output_path}")
            except Exception as e:
                self.status_label.text = f"Ошибка: {str(e)}"
                self.show_popup("Ошибка", f"Не удалось создать файл: {str(e)}")
        else:
            # Генерация нескольких бланков
            try:
                count = int(self.count_input.text)
                if count <= 0 or count > 1000:
                    self.show_popup("Ошибка", "Количество бланков должно быть от 1 до 1000!")
                    return
            except ValueError:
                self.show_popup("Ошибка", "Количество бланков должно быть числом!")
                return
            
            success_count, errors = generate_multiple_forms(template_path, form_number, count)
            
            if success_count > 0:
                status_message = f"Создано {success_count} из {count} маршрутных карт в папке {output_dir}"
                self.status_label.text = status_message
                
                if errors:
                    error_message = status_message + "\n\nОшибки:\n" + "\n".join(errors)
                    self.show_popup("Частичный успех", error_message)
                else:
                    self.show_popup("Успех", status_message)
            else:
                self.status_label.text = "Не удалось создать ни одной маршрутной карты"
                error_message = "Ошибки:\n" + "\n".join(errors)
                self.show_popup("Ошибка", error_message)

class FormGeneratorApp(App):
    def build(self):
        # Устанавливаем размер окна
        Window.size = (600, 500)
        # Устанавливаем минимальные размеры окна
        Window.minimum_width = 500
        Window.minimum_height = 400
        return FormGeneratorUI()

def update_empty_to_null():
    """Обновление пустых значений на NULL в существующей базе данных"""
    try:
        with sqlite3.connect('маршрутные_карты.db') as conn:
            cursor = conn.cursor()
            # Выполняем каждое обновление отдельно
            cursor.execute("UPDATE маршрутные_карты SET Учетный_номер = NULL WHERE Учетный_номер = ''")
            cursor.execute("UPDATE маршрутные_карты SET Номер_кластера = NULL WHERE Номер_кластера = ''")
            cursor.execute("UPDATE маршрутные_карты SET Статус = NULL WHERE Статус = ''")
            conn.commit()
    except sqlite3.Error as e:
        print(f"Ошибка при обновлении базы данных: {e}")
        raise

def main():
    update_empty_to_null()  # Добавьте эту строку перед запуском приложения
    FormGeneratorApp().run()

if __name__ == "__main__":
    main() 